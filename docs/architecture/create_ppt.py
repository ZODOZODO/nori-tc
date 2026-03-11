"""
nori-tc 차세대 아키텍처 PPT 생성 스크립트
템플릿: 0-DS-NoriMES Project 개요.pptx
출력: docs/architecture/nori-tc.pptx
"""
import sys
import copy
sys.path.insert(0, '/home/donggeon/.local/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ─── 경로 설정 ───
TEMPLATE_PATH = "/mnt/c/workspace_nori/work_2026/00.초기자료/0-DS-NoriMES Project 개요.pptx"
OUTPUT_PATH = "/home/donggeon/workspace_nori/nori-tc/docs/architecture/nori-tc.pptx"

# ─── 색상 팔레트 (기업 제안서 스타일) ───
COLOR_NAVY       = RGBColor(0x1A, 0x37, 0x5E)   # 진한 남색 (타이틀 배경)
COLOR_BLUE       = RGBColor(0x1F, 0x5C, 0x99)   # 중간 파란색 (강조)
COLOR_LIGHT_BLUE = RGBColor(0x2E, 0x86, 0xC1)   # 밝은 파란색 (부제 배경)
COLOR_ACCENT     = RGBColor(0xE8, 0x7C, 0x0F)   # 주황색 (포인트)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_GRAY_BG    = RGBColor(0xF4, 0xF6, 0xF9)
COLOR_GRAY_LINE  = RGBColor(0xCC, 0xCC, 0xCC)
COLOR_GREEN      = RGBColor(0x1A, 0x7A, 0x4A)   # 현재 기반 (초록)
COLOR_ORANGE     = RGBColor(0xD4, 0x5C, 0x0A)   # next 단계 (주황)
COLOR_PURPLE     = RGBColor(0x5A, 0x3E, 0x8E)   # AI 활용 (보라)

# ─── 슬라이드 크기 (16:9, 단위: EMU) ───
SLIDE_W = Cm(33.87)   # 12192000 EMU ≈ 33.87cm
SLIDE_H = Cm(19.05)   # 6858000 EMU  ≈ 19.05cm

# ─── 폰트 설정 ───
FONT_TITLE   = "맑은 고딕"
FONT_BODY    = "맑은 고딕"

# ─── 헬퍼 함수 ───

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Cm(left), Cm(top), Cm(width), Cm(height)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        if line_width:
            line.width = Pt(line_width)
    else:
        line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_name=FONT_BODY,
                 font_size=12, bold=False, color=COLOR_DARK_TEXT,
                 align=PP_ALIGN.LEFT, wrap=True):
    """텍스트박스 추가"""
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_text_box_multiline(slide, left, top, width, height, lines,
                           font_name=FONT_BODY, font_size=11,
                           color=COLOR_DARK_TEXT, line_spacing=None):
    """여러 줄 텍스트박스 추가"""
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, size, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size if size else font_size)
        run.font.bold = bold
        run.font.color.rgb = clr if clr else color
    return txBox


def add_slide_header(slide, title, subtitle=None, page_num=None):
    """슬라이드 공통 헤더 추가 (상단 네이비 바)"""
    # 상단 배경 바
    add_rect(slide, 0, 0, 33.87, 2.5, COLOR_NAVY)
    # 좌측 강조 선
    add_rect(slide, 0, 0, 0.5, 2.5, COLOR_ACCENT)
    # 제목
    add_text_box(slide, 0.8, 0.3, 28, 1.1, title,
                 font_name=FONT_TITLE, font_size=20, bold=True,
                 color=COLOR_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, 0.8, 1.4, 28, 0.9, subtitle,
                     font_name=FONT_BODY, font_size=12,
                     color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.LEFT)
    if page_num:
        add_text_box(slide, 31, 0.6, 2.5, 1.2, str(page_num),
                     font_size=14, bold=True, color=COLOR_WHITE,
                     align=PP_ALIGN.CENTER)
    # 하단 라인
    add_rect(slide, 0, 18.7, 33.87, 0.35, COLOR_ACCENT)


def add_keyword_badge(slide, left, top, width, text, bg_color, text_color=COLOR_WHITE):
    """키워드 뱃지 (라운드 느낌의 작은 사각형)"""
    add_rect(slide, left, top, width, 0.7, bg_color)
    add_text_box(slide, left + 0.1, top + 0.05, width - 0.2, 0.6, text,
                 font_size=10, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def add_section_box(slide, left, top, width, height, title, body_lines,
                    title_bg=COLOR_BLUE, body_bg=COLOR_GRAY_BG):
    """섹션 박스 (제목 배경 + 내용 배경)"""
    # 제목 영역
    add_rect(slide, left, top, width, 0.75, title_bg)
    add_text_box(slide, left + 0.2, top + 0.1, width - 0.4, 0.6, title,
                 font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)
    # 내용 영역
    add_rect(slide, left, top + 0.75, width, height - 0.75, body_bg,
             line_color=COLOR_GRAY_LINE, line_width=0.5)
    # 내용 텍스트
    txBox = slide.shapes.add_textbox(
        Cm(left + 0.2), Cm(top + 0.85), Cm(width - 0.4), Cm(height - 1.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.color.rgb = COLOR_DARK_TEXT


# ─── Presentation 생성 ───
prs = Presentation(TEMPLATE_PATH)

# 슬라이드 크기는 템플릿에서 그대로 가져옴 (12192000 x 6858000 EMU)
# 기존 슬라이드 전부 제거
R_ID_ATTR = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
slide_ids = [s for s in prs.slides._sldIdLst]
for sld_id in slide_ids:
    rId = sld_id.get(R_ID_ATTR)
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sld_id)


def add_new_slide(layout_idx=11):
    """새 슬라이드 추가 (Blank Slide 레이아웃 사용)"""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # 레이아웃에서 상속된 placeholder 제거
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


# ══════════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()

# 배경 전체 navy
add_rect(slide, 0, 0, 33.87, 19.05, COLOR_NAVY)

# 좌측 세로 강조 바
add_rect(slide, 0, 0, 0.6, 19.05, COLOR_ACCENT)

# 상단 장식 라인
add_rect(slide, 0.6, 0, 33.27, 0.15, RGBColor(0x2E, 0x86, 0xC1))

# 중앙 콘텐츠 영역 (반투명 효과 대신 약간 밝은 사각형)
add_rect(slide, 2.5, 4.5, 28.87, 10, RGBColor(0x1E, 0x44, 0x72))

# 제품명 태그
add_rect(slide, 2.5, 4.8, 7, 0.8, COLOR_ACCENT)
add_text_box(slide, 2.5, 4.85, 7, 0.7, "NEXT ARCHITECTURE 2025",
             font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# 메인 타이틀
add_text_box(slide, 2.5, 6.0, 28.87, 2.5, "nori-tc",
             font_name=FONT_TITLE, font_size=52, bold=True,
             color=COLOR_WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide, 2.5, 8.0, 28.87, 1.5, "라인 실행 플랫폼 아키텍처",
             font_name=FONT_TITLE, font_size=28, bold=False,
             color=RGBColor(0x90, 0xBE, 0xE8), align=PP_ALIGN.LEFT)

# 부제 설명
add_text_box(slide, 2.5, 9.8, 28.87, 1.2,
             "신규 설비 온보딩 · 변경 관리 · 메시지 신뢰성 · AI 운영 지원 · EDA/FDC 확장",
             font_size=13, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.LEFT)

# 가로 구분선
add_rect(slide, 2.5, 11.2, 28.87, 0.06, RGBColor(0x2E, 0x86, 0xC1))

# 날짜 / 버전
add_text_box(slide, 2.5, 11.5, 14, 0.9, "2025년  |  Proposal Version",
             font_size=11, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.LEFT)

# 6개 키워드 하단 뱃지
kw_colors = [COLOR_ACCENT, COLOR_BLUE, COLOR_LIGHT_BLUE,
             COLOR_GREEN, COLOR_PURPLE, COLOR_ORANGE]
kw_texts = ["Spec-to-Proof", "Verified ChangeOps", "Provable Reliability",
            "Explainable Ops", "Industrial Data Spine", "Commissioning Lab"]
kw_x = [2.5, 8.1, 13.7, 19.3, 24.2, 28.9]
kw_widths = [5.4, 5.4, 5.4, 4.7, 4.5, 4.8]
for i, (text, color, x, w) in enumerate(zip(kw_texts, kw_colors, kw_x, kw_widths)):
    add_rect(slide, x, 13.0, w, 0.75, color)
    add_text_box(slide, x + 0.1, 13.05, w - 0.2, 0.65, text,
                 font_size=9, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# 하단 장식
add_rect(slide, 0.6, 18.7, 33.27, 0.35, COLOR_ACCENT)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 2 — 목차
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "목차", "Table of Contents", 2)

# 배경
add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 목차 항목 데이터
toc_items = [
    ("1", "nori-tc 한 줄 정의 및 시장 배경",           "라인 실행 플랫폼으로서의 포지셔닝"),
    ("2", "왜 지금 nori-tc 인가",                      "고객 비용 구조와 6개 우선순위 키워드"),
    ("3", "Spec-to-Proof",                             "원본 문서 → evidence → 명세초안 → 승인"),
    ("4", "Verified ChangeOps",                        "변경 전 검증 · 반영 중 통제 · 복귀 절차"),
    ("5", "Provable Reliability",                      "파티션 라우팅 · DLQ · 발송이력 · 처리 근거"),
    ("6", "Explainable Operation Intelligence",        "사건 단위 타임라인 · AI 보조 운영 설명"),
    ("7", "Industrial Data Spine",                     "EDA/FDC 확장을 위한 기준 데이터 축"),
    ("8", "Autonomous Commissioning Lab",              "실라인 반영 전 동일 엔진 사전 검증"),
    ("9", "시장 경쟁사 비교 및 선택 이유",              "4가지 유형 14개 포지션 비교 분석"),
    ("10", "도입 기대 효과 및 결론",                   "비용 절감 포인트 · 도입 절차 · 기술 실사"),
]

num_colors = [COLOR_ACCENT, COLOR_ACCENT, COLOR_BLUE, COLOR_BLUE, COLOR_BLUE,
              COLOR_BLUE, COLOR_GREEN, COLOR_ORANGE, COLOR_NAVY, COLOR_NAVY]

col1_items = toc_items[:5]
col2_items = toc_items[5:]

for col, items in enumerate([col1_items, col2_items]):
    base_x = 1.2 + col * 16.8
    for row, (num, title, desc) in enumerate(items):
        top = 3.1 + row * 2.8
        # 번호 배지
        add_rect(slide, base_x, top, 1.2, 1.2, num_colors[row + col * 5])
        add_text_box(slide, base_x, top + 0.15, 1.2, 0.9, num,
                     font_size=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        # 제목 배경
        add_rect(slide, base_x + 1.2, top, 14.0 if col == 0 else 14.6, 0.75,
                 COLOR_NAVY if num_colors[row + col * 5] == COLOR_NAVY else COLOR_BLUE)
        add_text_box(slide, base_x + 1.4, top + 0.1, 13.5 if col == 0 else 14.0, 0.6,
                     title, font_size=11, bold=True, color=COLOR_WHITE)
        # 설명
        add_rect(slide, base_x + 1.2, top + 0.75, 14.0 if col == 0 else 14.6, 0.45,
                 COLOR_WHITE, line_color=COLOR_GRAY_LINE, line_width=0.3)
        add_text_box(slide, base_x + 1.4, top + 0.78, 13.5 if col == 0 else 14.0, 0.4,
                     desc, font_size=9, color=RGBColor(0x44, 0x55, 0x66))


# ══════════════════════════════════════════════════════════════════
# 슬라이드 3 — nori-tc 한 줄 정의 & 배경
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "nori-tc 한 줄 정의 및 시장 배경", "왜 지금 nori-tc인가", 3)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 정의 강조 박스
add_rect(slide, 1.0, 2.8, 31.87, 2.8, COLOR_NAVY)
add_rect(slide, 1.0, 2.8, 0.4, 2.8, COLOR_ACCENT)
add_text_box(slide, 1.6, 3.0, 30.87, 0.8,
             "nori-tc 한 줄 정의",
             font_size=13, bold=True, color=RGBColor(0xBB, 0xCC, 0xDD))
add_text_box(slide, 1.6, 3.7, 30.87, 1.6,
             "nori-tc는 설비를 단순 연결하는 TC가 아니라, 신규 설비 도입 속도·운영 변경 안정성·\n"
             "메시지 처리 증명 가능성·운영 설명 가능성·EDA/FDC 확장성을 함께 고려한 라인 실행 플랫폼입니다.",
             font_size=13, bold=False, color=COLOR_WHITE)

# 5가지 고객 비용 구조
add_text_box(slide, 1.0, 6.1, 20, 0.8, "고객이 실제로 비용을 지출하는 구간",
             font_size=14, bold=True, color=COLOR_NAVY)

cost_items = [
    ("신규 설비 도입 지연",     "생산 일정 지연으로 인한 비용"),
    ("변경 배포 중 라인 불안정", "생산 손실 발생 비용"),
    ("장애 원인 파악 지연",     "MTTR 증가로 인한 비용"),
    ("메시지 유실·중복 증명 불가", "운영 리스크 비용"),
    ("EDA/FDC 연동 시 구조 재설계", "고도화 재개발 비용"),
]

for i, (title, desc) in enumerate(cost_items):
    top = 7.1 + i * 1.8
    add_rect(slide, 1.0, top, 0.5, 1.3, COLOR_ACCENT)
    add_text_box(slide, 1.0, top + 0.3, 0.5, 0.7, str(i + 1),
                 font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.5, top, 18.5, 1.3, COLOR_WHITE, line_color=COLOR_GRAY_LINE, line_width=0.3)
    add_text_box(slide, 1.7, top + 0.1, 18.0, 0.7, title,
                 font_size=12, bold=True, color=COLOR_NAVY)
    add_text_box(slide, 1.7, top + 0.7, 18.0, 0.5, desc,
                 font_size=10, color=RGBColor(0x55, 0x66, 0x77))

# 우측 - 6개 키워드 박스
add_text_box(slide, 21.5, 6.1, 11.5, 0.8, "6개 핵심 키워드",
             font_size=14, bold=True, color=COLOR_NAVY)

kw_data = [
    ("1", "Spec-to-Proof",                    "신규 설비 온보딩 속도 향상",       COLOR_ACCENT),
    ("2", "Verified ChangeOps",               "변경 실패 리스크 감소",            COLOR_BLUE),
    ("3", "Provable Reliability",             "메시지 처리 결과 증명",            COLOR_LIGHT_BLUE),
    ("4", "Explainable Op Intelligence",      "장애 원인 설명·대응 속도 향상",    COLOR_GREEN),
    ("5", "Industrial Data Spine",            "EDA/FDC 고도화 준비",             COLOR_PURPLE),
    ("6", "Autonomous Commissioning Lab",     "실라인 반영 전 사전 검증",         COLOR_ORANGE),
]

for i, (num, kw, desc, color) in enumerate(kw_data):
    top = 7.1 + i * 1.8
    add_rect(slide, 21.5, top, 11.5, 1.3, COLOR_WHITE, line_color=color, line_width=1.5)
    add_rect(slide, 21.5, top, 1.2, 1.3, color)
    add_text_box(slide, 21.5, top + 0.3, 1.2, 0.7, num,
                 font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 22.8, top + 0.05, 9.9, 0.65, kw,
                 font_size=11, bold=True, color=color)
    add_text_box(slide, 22.8, top + 0.65, 9.9, 0.55, desc,
                 font_size=9, color=RGBColor(0x44, 0x55, 0x66))


# ══════════════════════════════════════════════════════════════════
# 슬라이드 4 — Spec-to-Proof
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "Spec-to-Proof",
                 "원본 설비 문서 → evidence 정규화 → 명세 초안 → 사람 승인 흐름", 4)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 고객 문제
add_rect(slide, 0.8, 2.8, 32.0, 1.6, RGBColor(0xFF, 0xF3, 0xE0))
add_rect(slide, 0.8, 2.8, 0.35, 1.6, COLOR_ACCENT)
add_text_box(slide, 1.3, 2.9, 31, 0.6, "고객 문제",
             font_size=11, bold=True, color=COLOR_ACCENT)
add_text_box(slide, 1.3, 3.4, 31, 0.8,
             "설비사 문서 해석, 메시지 규칙 정리, parser 작성, workflow 초안 작성이 개발자 개인 역량에 의존 → "
             "신규 설비 도입 4~8주 이상 소요. 리드타임의 30~50% 단축 가능.",
             font_size=10.5, color=COLOR_DARK_TEXT)

# 처리 흐름 (6단계 화살표)
flow_steps = [
    ("① 문서 입력",      "PDF/Word\nExcel/XML\n샘플 로그"),
    ("② 형식별 추출",    "규칙 기반 파서\n문서 분석 엔진\nLog 분해"),
    ("③ evidence\n정규화", "공통 구조\n원본 위치\n신뢰도 포함"),
    ("④ RAG/LLM\n초안 생성", "메시지 규칙\nworkflow\naction 초안"),
    ("⑤ AI Agent\n교차 검토", "누락 필드\n상충 규칙\n확인 질문"),
    ("⑥ 사람 승인\n→ 반영", "Parser 생성\nworkflow 반영\n테스트 준비"),
]

step_colors = [COLOR_NAVY, COLOR_BLUE, COLOR_LIGHT_BLUE,
               COLOR_GREEN, COLOR_PURPLE, COLOR_ACCENT]
step_w = 4.8
for i, (title, desc) in enumerate(flow_steps):
    x = 0.8 + i * (step_w + 0.5)
    add_rect(slide, x, 4.8, step_w, 1.0, step_colors[i])
    add_text_box(slide, x + 0.1, 4.85, step_w - 0.2, 0.9, title,
                 font_size=10, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 5.8, step_w, 1.8, COLOR_WHITE,
             line_color=step_colors[i], line_width=1.0)
    add_text_box(slide, x + 0.15, 5.9, step_w - 0.3, 1.6, desc,
                 font_size=9.5, color=COLOR_DARK_TEXT, align=PP_ALIGN.CENTER)
    # 화살표 (마지막 제외)
    if i < len(flow_steps) - 1:
        add_text_box(slide, x + step_w + 0.05, 5.0, 0.5, 0.7, "→",
                     font_size=18, bold=True, color=COLOR_NAVY, align=PP_ALIGN.CENTER)

# 현재 기반 / next 단계 / 비용절감
sections = [
    ("현재 기반", COLOR_GREEN,
     ["• 모델 스키마 정의 완료 (메시지 타입·파라미터·workflow·MDF 구조)",
      "• 모델 버전 관리 구조",
      "• 조회 및 반영 API / UI 구조 존재"]),
    ("next 단계에서 추가", COLOR_ORANGE,
     ["• 문서 입력/추출 계층 및 형식별 추출 어댑터",
      "• evidence block 표준 스키마",
      "• RAG 인덱싱 계층 (모델 스키마 + 온보딩 사례)",
      "• 명세 초안 검토/승인 UI"]),
    ("비용 절감 포인트", COLOR_ACCENT,
     ["• 설비 도입 문서 해석·초안 작성 30~50% 단축",
      "• 벤더 재질의 / 재작업 감소",
      "• 숙련 개발자 투입 시간 감소"]),
]

for i, (title, color, items) in enumerate(sections):
    x = 0.8 + i * 11.2
    add_rect(slide, x, 8.2, 10.8, 0.75, color)
    add_text_box(slide, x + 0.2, 8.25, 10.4, 0.65, title,
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, 8.95, 10.8, 9.4, COLOR_WHITE,
             line_color=color, line_width=1.0)
    for j, item in enumerate(items):
        add_text_box(slide, x + 0.2, 9.1 + j * 1.0, 10.4, 0.9, item,
                     font_size=10, color=COLOR_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 5 — Verified ChangeOps
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "Verified ChangeOps",
                 "변경 전 검증 · 반영 중 통제 · 반영 후 복구를 하나의 운영 절차로", 5)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 고객 문제
add_rect(slide, 0.8, 2.8, 32.0, 1.6, RGBColor(0xFF, 0xF3, 0xE0))
add_rect(slide, 0.8, 2.8, 0.35, 1.6, COLOR_ACCENT)
add_text_box(slide, 1.3, 2.9, 31, 0.5, "고객 문제",
             font_size=11, bold=True, color=COLOR_ACCENT)
add_text_box(slide, 1.3, 3.4, 31, 0.8,
             "운영 중 변경 실패가 가장 큰 리스크. 계획된 변경 작업의 15~25%에서 예기치 않은 문제 발생. "
             "라인 중단 1시간당 수천만~수억 원의 손실.",
             font_size=10.5, color=COLOR_DARK_TEXT)

# 6단계 변경 절차
steps = [
    ("운영 버전\n분리",     "현재 운영 버전과\n편집 버전 분리"),
    ("Diff\n시각화",        "사람이 읽을 수 있는\nDiff 표시"),
    ("영향도\n계산",        "EQP·model·workflow·\ngateway 영향 범위"),
    ("사전\n검증",          "운영 로그 기반\n사전 검증 수행"),
    ("단계\n반영",          "일부 설비/gateway group\n부터 순차 반영"),
    ("복귀\n절차",          "이상 징후 감지 시\n이전 버전으로 복귀"),
]

colors_steps = [COLOR_NAVY, COLOR_BLUE, COLOR_LIGHT_BLUE,
                COLOR_GREEN, COLOR_ACCENT, COLOR_ORANGE]

for i, (title, desc) in enumerate(steps):
    x = 0.8 + i * 5.5
    # 번호 원
    add_rect(slide, x + 1.2, 4.8, 1.4, 1.4, colors_steps[i])
    add_text_box(slide, x + 1.2, 4.95, 1.4, 0.8, str(i + 1),
                 font_size=20, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # 제목
    add_rect(slide, x, 6.3, 4.5, 0.85, colors_steps[i])
    add_text_box(slide, x + 0.1, 6.35, 4.3, 0.75, title.replace('\n', ' '),
                 font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # 설명
    add_rect(slide, x, 7.15, 4.5, 1.6, COLOR_WHITE,
             line_color=colors_steps[i], line_width=1.0)
    add_text_box(slide, x + 0.15, 7.25, 4.2, 1.4, desc,
                 font_size=10, color=COLOR_DARK_TEXT, align=PP_ALIGN.CENTER)
    # 화살표
    if i < len(steps) - 1:
        add_text_box(slide, x + 4.5, 5.2, 0.5, 0.7, "→",
                     font_size=16, bold=True, color=COLOR_NAVY, align=PP_ALIGN.CENTER)

# AI 역할 구분
add_rect(slide, 0.8, 9.3, 32.0, 0.7, COLOR_PURPLE)
add_text_box(slide, 1.0, 9.35, 31.5, 0.6,
             "AI 보조 역할: diff 요약 · 영향도 설명 · 검토 포인트 추천 · 승인 문구/결과 요약 초안 (실제 반영 통제는 규칙 기반 메커니즘)",
             font_size=10.5, bold=True, color=COLOR_WHITE)

# 현재 기반 / next 단계 / 비용절감
sections = [
    ("현재 기반", COLOR_GREEN,
     ["• Check Out → EDIT → Save → Check In 구조",
      "• 모델 EDIT 버전 변경 흐름",
      "• 고정 partition 기반 gateway 라우팅 구조"]),
    ("next 단계에서 추가", COLOR_ORANGE,
     ["• 변경 diff 표준화",
      "• 영향도 계산 그래프",
      "• 운영 로그 기반 사전 검증",
      "• 단계 반영 계획 UI",
      "• 이상 징후 감지와 복귀 절차 자동화"]),
    ("비용 절감 포인트", COLOR_ACCENT,
     ["• 변경 사고로 인한 라인 중단 비용 감소",
      "• 긴급 롤백 대응 비용 감소 (60~80% 단축 목표)",
      "• 변경 검토/승인/운영 이관 시간 감소"]),
]

for i, (title, color, items) in enumerate(sections):
    x = 0.8 + i * 11.2
    add_rect(slide, x, 10.3, 10.8, 0.75, color)
    add_text_box(slide, x + 0.2, 10.35, 10.4, 0.65, title,
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, 11.05, 10.8, 7.5, COLOR_WHITE,
             line_color=color, line_width=1.0)
    for j, item in enumerate(items):
        add_text_box(slide, x + 0.2, 11.2 + j * 1.0, 10.4, 0.9, item,
                     font_size=10, color=COLOR_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 6 — Provable Reliability
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "Provable Reliability",
                 "파티션 라우팅 · Dual Response · DLQ · 발송 이력 구조로 처리 결과를 증명", 6)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 고객 문제
add_rect(slide, 0.8, 2.8, 32.0, 1.4, RGBColor(0xFF, 0xF3, 0xE0))
add_rect(slide, 0.8, 2.8, 0.35, 1.4, COLOR_ACCENT)
add_text_box(slide, 1.3, 2.9, 31, 0.5, "고객 문제",
             font_size=11, bold=True, color=COLOR_ACCENT)
add_text_box(slide, 1.3, 3.3, 31, 0.7,
             "단순히 '안정적'이라는 표현이 아닌, 메시지 생성 → 파티션 발행 → 처리 → 실패 이력까지 추적·설명 가능해야 함.",
             font_size=10.5, color=COLOR_DARK_TEXT)

# 핵심 메커니즘 6가지
mechs = [
    ("route_partition\n기반 라우팅",   "tc_eqp.route_partition을\n단일 라우팅 기준으로\n명시 partition 발행"),
    ("Split-brain\n방지",              "같은 PASSIVE listener-group\n→ 같은 partition 강제"),
    ("Dual Response\n수집",            "EQP_CREATE/UPDATE/DELETE\n→ Gateway + Business\n양쪽 응답 모두 완료"),
    ("Async\nTrace 저장",              "EQP_START/END\n비동기 trace 단위 추적\n결과 별도 저장"),
    ("DLQ\n격리",                      "처리 실패 시\nGateway/Business DLQ\n격리 → 재처리 가능"),
    ("Transactional\nOutbox",          "outbox 테이블 설계 완료\nnext 단계에서\nend-to-end 연결"),
]

colors_mechs = [COLOR_NAVY, COLOR_BLUE, COLOR_LIGHT_BLUE,
                COLOR_GREEN, COLOR_ORANGE, COLOR_PURPLE]

for i, (title, desc) in enumerate(mechs):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 11.1
    top = 4.6 + row * 4.0
    add_rect(slide, x, top, 10.7, 0.85, colors_mechs[i])
    add_text_box(slide, x + 0.2, top + 0.1, 10.3, 0.7, title.replace('\n', ' '),
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, top + 0.85, 10.7, 2.6, COLOR_WHITE,
             line_color=colors_mechs[i], line_width=1.0)
    add_text_box(slide, x + 0.2, top + 1.0, 10.3, 2.3, desc,
                 font_size=10.5, color=COLOR_DARK_TEXT)

# 현재 기반 / next / 비용절감
sections = [
    ("현재 기반", COLOR_GREEN,
     ["• 고정 partition 라우팅 규칙",
      "• Dual Response 수집 구조",
      "• Async 결과 저장 구조",
      "• Gateway / Business DLQ",
      "• Transactional Outbox 테이블 설계"]),
    ("next 단계에서 추가", COLOR_ORANGE,
     ["• 런타임 발행 경로와 outbox end-to-end 연결",
      "• 재발행 워커",
      "• 발송 이력 / 재시도 / dead 상태 운영 뷰",
      "• 신뢰성 진단 대시보드"]),
    ("비용 절감 포인트", COLOR_ACCENT,
     ["• 장애 재현 인건비 감소 (trace 기반 조회)",
      "• 유실/중복 의심 건 조사 시간 감소",
      "• 현장 신뢰 이슈 운영 리스크 감소"]),
]

for i, (title, color, items) in enumerate(sections):
    x = 0.8 + i * 11.2
    add_rect(slide, x, 13.0, 10.8, 0.75, color)
    add_text_box(slide, x + 0.2, 13.05, 10.4, 0.65, title,
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, 13.75, 10.8, 4.8, COLOR_WHITE,
             line_color=color, line_width=1.0)
    for j, item in enumerate(items):
        add_text_box(slide, x + 0.2, 13.9 + j * 0.9, 10.4, 0.8, item,
                     font_size=10, color=COLOR_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 7 — Explainable Operation Intelligence
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "Explainable Operation Intelligence",
                 "사건 단위 타임라인 + AI 보조로 운영자가 추적해야 할 근거를 더 빨리 묶어줌", 7)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 고객 문제
add_rect(slide, 0.8, 2.8, 32.0, 1.4, RGBColor(0xFF, 0xF3, 0xE0))
add_rect(slide, 0.8, 2.8, 0.35, 1.4, COLOR_ACCENT)
add_text_box(slide, 1.3, 2.9, 31, 0.5, "고객 문제",
             font_size=11, bold=True, color=COLOR_ACCENT)
add_text_box(slide, 1.3, 3.3, 31, 0.7,
             "장애 원인 파악까지 평균 2~6시간, 대응 완료까지 4~10시간. 숙련 운영자만 빠르게 복원 가능 → 야간/교대 시 더 길어짐.",
             font_size=10.5, color=COLOR_DARK_TEXT)

# 동작 방향 (좌측 흐름)
add_text_box(slide, 0.8, 4.5, 19, 0.75, "동작 방식",
             font_size=13, bold=True, color=COLOR_NAVY)

flow_items = [
    ("1", "데이터 수집", "TC 로그, traceId, eqpId, modelVersion, workflow,\nruntime state, DLQ, 작업 이력을 사건 단위로 묶음"),
    ("2", "사건 타임라인 구성", "어떤 설비 → 어떤 변경 이후 → 어떤 메시지 실패 →\n상태 전이 → 현재 영향 범위 → 다음 대응"),
    ("3", "RAG 조회", "관련 로그, 상태, 버전 이력, 작업 문맥 함께 조회"),
    ("4", "LLM 설명 생성", "사건 설명, 가능한 원인, 다음 조치 정리"),
    ("5", "AI Agent 검토", "추가 조회 반복, 근거 부족 시 '추정'으로 명시"),
    ("6", "근거 링크 표시", "어떤 로그/이벤트가 근거인지 화면에 함께 표시"),
]

for i, (num, title, desc) in enumerate(flow_items):
    top = 5.4 + i * 2.1
    add_rect(slide, 0.8, top, 0.8, 1.6, COLOR_NAVY)
    add_text_box(slide, 0.8, top + 0.35, 0.8, 0.8, num,
                 font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.6, top, 18.0, 0.7, COLOR_BLUE)
    add_text_box(slide, 1.7, top + 0.08, 17.7, 0.58, title,
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_rect(slide, 1.6, top + 0.7, 18.0, 0.9, COLOR_WHITE,
             line_color=COLOR_BLUE, line_width=0.5)
    add_text_box(slide, 1.8, top + 0.75, 17.5, 0.85, desc,
                 font_size=9.5, color=COLOR_DARK_TEXT)

# 우측 - 현재기반/next/비용절감
right_x = 21.0
sections = [
    ("현재 기반", COLOR_GREEN,
     ["• traceId 기반 요청/응답 추적 구조",
      "• EQP runtime-state 조회 구조",
      "• Work→ControlJob→ProcessJob→Lot 계층",
      "• 모델/워크플로/파라미터 버전 구조"]),
    ("next 단계에서 추가", COLOR_ORANGE,
     ["• 사건 단위 타임라인 뷰",
      "• 근거 링크가 포함된 자연어 설명",
      "• 대응 체크리스트 생성",
      "• 교대 근무 인수인계용 사건 요약"]),
    ("비용 절감 포인트", COLOR_ACCENT,
     ["• MTTR 50~70% 단축 목표",
      "• 고급 운영 인력 의존도 감소",
      "• 장애 리포트 작성 시간 감소"]),
]

for i, (title, color, items) in enumerate(sections):
    top = 4.5 + i * 4.8
    add_rect(slide, right_x, top, 12.5, 0.75, color)
    add_text_box(slide, right_x + 0.2, top + 0.1, 12.1, 0.65, title,
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_rect(slide, right_x, top + 0.75, 12.5, 3.8, COLOR_WHITE,
             line_color=color, line_width=1.0)
    for j, item in enumerate(items):
        add_text_box(slide, right_x + 0.2, top + 0.95 + j * 0.88, 12.1, 0.8, item,
                     font_size=10, color=COLOR_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 8 — Industrial Data Spine & Autonomous Commissioning Lab
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "Industrial Data Spine  +  Autonomous Commissioning Lab",
                 "EDA/FDC 확장을 위한 기준 데이터 축 · 실라인 반영 전 동일 엔진 사전 검증", 8)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 가운데 구분선
add_rect(slide, 16.7, 2.7, 0.15, 15.7, COLOR_GRAY_LINE)

# ── 좌측: Industrial Data Spine ──
add_rect(slide, 0.8, 2.8, 15.6, 0.85, COLOR_GREEN)
add_text_box(slide, 1.0, 2.88, 15.2, 0.7, "Industrial Data Spine",
             font_size=13, bold=True, color=COLOR_WHITE)

add_text_box(slide, 0.8, 3.85, 15.6, 1.0,
             "TC를 단일 목적 시스템으로 끝내지 않고,\nEDA/FDC/AI가 공통으로 기대는 기준 데이터 축 구성",
             font_size=10.5, color=COLOR_DARK_TEXT)

# 기준 축
spine_items = ["EQP (설비)", "Model / Model Version", "Workflow (공정 흐름)",
               "Runtime State / Alarm", "Work / Carrier / Lot 계층"]
add_text_box(slide, 0.8, 5.0, 15, 0.7, "핵심 기준 데이터 축",
             font_size=11, bold=True, color=COLOR_GREEN)
for i, item in enumerate(spine_items):
    add_rect(slide, 0.8, 5.8 + i * 0.85, 15.6, 0.75,
             COLOR_WHITE, line_color=COLOR_GREEN, line_width=0.8)
    add_text_box(slide, 1.0, 5.87 + i * 0.85, 15.2, 0.65, f"• {item}",
                 font_size=10.5, color=COLOR_DARK_TEXT)

# EDA / FDC 방향
add_text_box(slide, 0.8, 10.2, 15, 0.7, "확장 방향",
             font_size=11, bold=True, color=COLOR_GREEN)

for i, (dir_name, dir_desc, color) in enumerate([
    ("EDA 방향", "고속 수집 데이터를 EQP·model version·\n상태·작업 문맥과 연결", COLOR_BLUE),
    ("FDC 방향", "recipe·chamber·lot·작업 상태와\n정렬된 분석 가능 데이터로 변환", COLOR_PURPLE),
]):
    x = 0.8 + i * 8.0
    add_rect(slide, x, 11.0, 7.6, 0.75, color)
    add_text_box(slide, x + 0.2, 11.08, 7.2, 0.65, dir_name,
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, 11.75, 7.6, 2.0, COLOR_WHITE,
             line_color=color, line_width=0.8)
    add_text_box(slide, x + 0.2, 11.9, 7.2, 1.7, dir_desc,
                 font_size=10, color=COLOR_DARK_TEXT)

# next / 비용절감
add_rect(slide, 0.8, 14.0, 7.6, 0.7, COLOR_ORANGE)
add_text_box(slide, 1.0, 14.06, 7.2, 0.6, "next 단계에서 추가",
             font_size=10, bold=True, color=COLOR_WHITE)
add_rect(slide, 0.8, 14.7, 7.6, 3.8, COLOR_WHITE,
         line_color=COLOR_ORANGE, line_width=0.8)
for j, item in enumerate(["• 공통 이벤트/데이터 계약 표준화", "• 분석용 파생 데이터셋",
                           "• EDA/FDC 연계용 기준 키 매핑", "• RAG 지식 베이스 운영 용어사전"]):
    add_text_box(slide, 1.0, 14.85 + j * 0.85, 7.2, 0.8, item,
                 font_size=10, color=COLOR_DARK_TEXT)

add_rect(slide, 8.6, 14.0, 7.8, 0.7, COLOR_ACCENT)
add_text_box(slide, 8.8, 14.06, 7.4, 0.6, "비용 절감 포인트",
             font_size=10, bold=True, color=COLOR_WHITE)
add_rect(slide, 8.6, 14.7, 7.8, 3.8, COLOR_WHITE,
         line_color=COLOR_ACCENT, line_width=0.8)
for j, item in enumerate(["• EDA/FDC 착수 시 데이터 재정의 최소화", "• 연동 불일치 감소",
                           "• 후속 프로젝트 기간 단축"]):
    add_text_box(slide, 8.8, 14.85 + j * 0.85, 7.4, 0.8, item,
                 font_size=10, color=COLOR_DARK_TEXT)

# ── 우측: Autonomous Commissioning Lab ──
rx = 17.1
add_rect(slide, rx, 2.8, 16.2, 0.85, COLOR_ORANGE)
add_text_box(slide, rx + 0.2, 2.88, 15.8, 0.7, "Autonomous Commissioning Lab",
             font_size=13, bold=True, color=COLOR_WHITE)

add_text_box(slide, rx, 3.85, 16.2, 1.0,
             "실라인 반영 전에 가능한 한 같은 엔진으로 먼저 검증하는 사전 검증 환경",
             font_size=10.5, color=COLOR_DARK_TEXT)

# 5단계 검증 흐름
lab_steps = [
    ("명세 입력", "신규/변경 명세 입력"),
    ("로그 준비", "샘플/운영 로그\n시뮬레이터 입력"),
    ("동일 엔진\n검증 실행", "gateway runtime\nworkflow 엔진\n상태 전이 규칙"),
    ("결과 분석", "parser/workflow\ntimeout/start/end\n흐름 검사"),
    ("결과서\n생성", "테스트 결과서\n자동 정리"),
]

lab_colors = [COLOR_NAVY, COLOR_BLUE, COLOR_ORANGE, COLOR_GREEN, COLOR_PURPLE]

for i, (title, desc) in enumerate(lab_steps):
    x = rx + 0.3 + i * 3.1
    add_rect(slide, x, 5.0, 2.7, 0.85, lab_colors[i])
    add_text_box(slide, x + 0.05, 5.05, 2.6, 0.75, title.replace('\n', ' '),
                 font_size=9.5, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 5.85, 2.7, 2.0, COLOR_WHITE,
             line_color=lab_colors[i], line_width=0.8)
    add_text_box(slide, x + 0.1, 5.95, 2.5, 1.8, desc,
                 font_size=9, color=COLOR_DARK_TEXT, align=PP_ALIGN.CENTER)
    if i < len(lab_steps) - 1:
        add_text_box(slide, x + 2.7, 5.35, 0.35, 0.7, "→",
                     font_size=14, bold=True, color=COLOR_NAVY, align=PP_ALIGN.CENTER)

# AI 보조 역할
add_rect(slide, rx, 8.2, 16.2, 0.75, COLOR_PURPLE)
add_text_box(slide, rx + 0.2, 8.26, 15.8, 0.65,
             "AI 보조: 테스트 시나리오 초안 생성 · 누락 케이스 추천 · 결과 요약 및 이상 포인트 설명",
             font_size=10.5, bold=True, color=COLOR_WHITE)

# 현재기반 / next / 비용절감
lab_sections = [
    ("현재 기반", COLOR_GREEN,
     ["• Gateway runtime control 구조",
      "• async lifecycle 처리 구조",
      "• 모델/워크플로/설비 상태 기반 구조"]),
    ("next 단계에서 추가", COLOR_ORANGE,
     ["• 샘플/운영 로그 replay 환경",
      "• 시뮬레이터 연계 어댑터",
      "• 명세 기반 테스트 실행기",
      "• 결과서 자동 정리"]),
    ("비용 절감 포인트", COLOR_ACCENT,
     ["• 현장 시험 반복 비용 감소",
      "• 초기 안정화 기간 단축",
      "• 설비 반영 후 긴급 대응 비용 감소"]),
]

for i, (title, color, items) in enumerate(lab_sections):
    top = 9.3 + i * 3.1
    add_rect(slide, rx, top, 16.2, 0.75, color)
    add_text_box(slide, rx + 0.2, top + 0.1, 15.8, 0.65, title,
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_rect(slide, rx, top + 0.75, 16.2, 2.1, COLOR_WHITE,
             line_color=color, line_width=0.8)
    for j, item in enumerate(items):
        add_text_box(slide, rx + 0.2, top + 0.9 + j * 0.62, 15.8, 0.6, item,
                     font_size=10, color=COLOR_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 9 — AI 인프라 구성 및 보안 원칙
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "AI 인프라 구성 및 보안 원칙",
                 "온프레미스 전용 AI · 데이터 외부 미전송 · LLM 구성 가이드", 9)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 데이터 보안 강조 박스
add_rect(slide, 0.8, 2.8, 32.0, 2.0, COLOR_NAVY)
add_rect(slide, 0.8, 2.8, 0.4, 2.0, COLOR_ACCENT)
add_text_box(slide, 1.4, 2.9, 30.9, 0.75, "데이터 보안 원칙",
             font_size=13, bold=True, color=COLOR_ACCENT)
add_text_box(slide, 1.4, 3.55, 30.9, 1.0,
             "모든 AI 처리(문서 추출, RAG 검색, LLM 추론, Agent 실행)는 고객사 내부 서버에서 온프레미스/폐쇄망 환경으로 수행. "
             "설비 명세·공정 로그·모델 데이터는 고객사 네트워크 외부로 나가지 않습니다.",
             font_size=11, color=COLOR_WHITE)

# GPU VRAM별 LLM 구성 표
add_text_box(slide, 0.8, 5.1, 20, 0.8, "GPU VRAM별 추천 LLM 구성",
             font_size=13, bold=True, color=COLOR_NAVY)

headers = ["GPU VRAM", "추천 모델", "특징"]
col_widths = [4.5, 6.0, 21.5]
row_data = [
    ("16GB 이상",    "Qwen2.5-7B",     "한국어/영어/중국어 다국어 지원, 구조화 출력 품질 양호"),
    ("24GB 이상",    "Qwen2.5-14B",    "7B 대비 명세 추출과 교차 검토 정확도 향상"),
    ("40GB 이상\n(단일/멀티 GPU)", "Qwen2.5-72B\n/ Llama 3.1-70B", "복잡한 문서 교차 분석, 사건 설명 품질 최상"),
]

# 헤더
x = 0.8
for j, (header, w) in enumerate(zip(headers, col_widths)):
    add_rect(slide, x, 6.1, w, 0.75, COLOR_NAVY)
    add_text_box(slide, x + 0.1, 6.15, w - 0.2, 0.65, header,
                 font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    x += w

row_bg_colors = [COLOR_WHITE, COLOR_GRAY_BG, COLOR_WHITE]
for i, (vram, model, desc) in enumerate(row_data):
    x = 0.8
    top = 6.85 + i * 1.15
    rh = 1.1
    for j, (text, w) in enumerate(zip([vram, model, desc], col_widths)):
        add_rect(slide, x, top, w, rh, row_bg_colors[i],
                 line_color=COLOR_GRAY_LINE, line_width=0.3)
        add_text_box(slide, x + 0.1, top + 0.05, w - 0.2, rh - 0.1, text,
                     font_size=10, color=COLOR_DARK_TEXT, align=PP_ALIGN.CENTER if j < 2 else PP_ALIGN.LEFT)
        x += w

# Qwen 추천 이유
add_rect(slide, 0.8, 10.4, 32.0, 0.7, RGBColor(0xE8, 0xF4, 0xE8))
add_rect(slide, 0.8, 10.4, 0.35, 0.7, COLOR_GREEN)
add_text_box(slide, 1.3, 10.45, 31.0, 0.6,
             "Qwen2.5 시리즈 1순위 추천 이유: 한국어와 제조 도메인 용어 이해 품질이 높고, "
             "JSON 등 구조화 형식 출력에 강점 → evidence block 생성과 명세 초안 생성에 적합",
             font_size=10.5, color=COLOR_DARK_TEXT)

# AI 오류 관리
add_text_box(slide, 0.8, 11.4, 32, 0.8, "AI 오류 관리 (Hallucination 대응)",
             font_size=13, bold=True, color=COLOR_NAVY)

ai_items = [
    ("사람 승인 게이트 필수", COLOR_ACCENT,
     "AI가 생성한 명세 초안과 운영 대응 제안은 사람이 검토·승인 후에만 실제 운영에 반영. AI 단독으로 런타임 반영 불가."),
    ("근거 링크 표시", COLOR_BLUE,
     "AI가 설명 생성 시 어떤 로그·이벤트·문서 원본을 근거로 했는지 함께 표시. 운영자가 직접 근거 확인 가능."),
    ("추정 명시", COLOR_PURPLE,
     "근거가 부족한 부분은 '추정'으로 명확히 표시하거나 운영자 확인 질문으로 전환. AI가 모르는 것을 아는 척하지 않음."),
]

for i, (title, color, desc) in enumerate(ai_items):
    x = 0.8 + i * 11.1
    add_rect(slide, x, 12.4, 10.7, 0.75, color)
    add_text_box(slide, x + 0.2, 12.45, 10.3, 0.65, title,
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, 13.15, 10.7, 5.2, COLOR_WHITE,
             line_color=color, line_width=1.0)
    add_text_box(slide, x + 0.2, 13.3, 10.3, 4.8, desc,
                 font_size=10.5, color=COLOR_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 10 — 일반 TC 대비 nori-tc 차이
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "일반적인 TC 대비 nori-tc의 차이",
                 "더 검증 가능하게 · 더 운영 가능하게 · 더 확장 가능하게", 10)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 비교표 헤더
headers = ["비교 항목", "일반적인 TC", "nori-tc"]
col_widths = [7.0, 11.0, 14.87]
header_colors = [COLOR_NAVY, RGBColor(0x6D, 0x6D, 0x6D), COLOR_BLUE]

x = 0.8
for header, w, hc in zip(headers, col_widths, header_colors):
    add_rect(slide, x, 2.8, w, 0.85, hc)
    add_text_box(slide, x + 0.15, 2.88, w - 0.3, 0.7, header,
                 font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    x += w

rows = [
    ("신규 설비 대응", "문서 해석 후 수작업 개발 중심",
     "원본 문서 → evidence 정규화 → 명세 초안 → 검토/승인 → 구현"),
    ("변경 관리", "운영 반영 전 수동 점검 중심",
     "버전 분리, diff, 영향도, 사전 검증, 단계 반영, 복귀 절차"),
    ("안정성 설명", "'안정적'이라는 표현 위주",
     "partition, 응답 수집, 비동기 trace, DLQ, 발송 이력 구조를 근거로 설명"),
    ("장애 대응", "로그 검색과 숙련 인력 의존",
     "trace, 상태, workflow, 작업 이력 기반 사건 단위 설명"),
    ("고도화 확장", "TC 자체에 머무름",
     "EDA/FDC/AI 활용까지 이어지는 기준 데이터 축 제공"),
]

row_height = 2.3
for i, (item, general, nori) in enumerate(rows):
    bg = COLOR_WHITE if i % 2 == 0 else COLOR_GRAY_BG
    x = 0.8
    top = 3.65 + i * row_height
    for text, w in zip([item, general, nori], col_widths):
        add_rect(slide, x, top, w, row_height, bg,
                 line_color=COLOR_GRAY_LINE, line_width=0.3)
        text_color = COLOR_NAVY if x == 0.8 else (
            RGBColor(0x55, 0x55, 0x55) if w == col_widths[1] else COLOR_DARK_TEXT)
        add_text_box(slide, x + 0.15, top + 0.15, w - 0.3, row_height - 0.3, text,
                     font_size=10.5, bold=(x == 0.8), color=text_color)
        x += w

# 하단 메시지
add_rect(slide, 0.8, 15.2, 32.0, 1.0, COLOR_NAVY)
add_text_box(slide, 1.1, 15.3, 31.4, 0.8,
             "nori-tc의 차별점은 같은 TC 기능을 '더 검증 가능하게', '더 운영 가능하게', '더 확장 가능하게' 만드는 방향에 있습니다.",
             font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 11 — 시장 경쟁사 비교 (글로벌/엔터프라이즈)
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "시장 경쟁사 비교 — 글로벌 및 엔터프라이즈",
                 "4가지 유형 14개 포지션에서 nori-tc의 위치", 11)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 4개 포지션 박스
positions = [
    ("설비 OEM 측 도구",
     "Cimetrix, PEER Group\n코닉오토메이션 EasyCluster",
     "설비 내부 SECS/GEM 연결 솔루션\n팹 운영자의 TC 역할 담당 안함\n→ nori-tc와 보완 관계",
     COLOR_GRAY_LINE),
    ("엔터프라이즈 MES 패키지",
     "IBM SiView\nSiemens Opcenter",
     "TC가 MES 전체 패키지 일부\nTC만 단독 도입 불가\n수억 원 이상 초기 투자 + 수년 구축",
     RGBColor(0xCC, 0x44, 0x44)),
    ("국내 통합 플랫폼",
     "에임시스템 Nexbe+ EAS",
     "20년+ 국내 레퍼런스\nMES 전체와 묶여 TC 단독 불가\n변경관리·AI 기능 부재",
     COLOR_ORANGE),
    ("자체 개발 TC",
     "삼성전자, SK하이닉스",
     "완전 맞춤 구현 가능\n수백 명 내부 개발 조직 필요\n중견 팹/신규 라인에는 비현실적",
     COLOR_PURPLE),
]

for i, (category, products, desc, color) in enumerate(positions):
    x = 0.8 + i * 8.1
    add_rect(slide, x, 2.8, 7.7, 0.75, color if color != COLOR_GRAY_LINE else COLOR_NAVY)
    add_text_box(slide, x + 0.15, 2.88, 7.4, 0.65, category,
                 font_size=10.5, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, 3.55, 7.7, 1.3, COLOR_WHITE,
             line_color=color if color != COLOR_GRAY_LINE else COLOR_NAVY, line_width=1.0)
    add_text_box(slide, x + 0.15, 3.65, 7.4, 1.1, products,
                 font_size=10, bold=True, color=COLOR_NAVY, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 4.85, 7.7, 2.8, COLOR_GRAY_BG,
             line_color=COLOR_GRAY_LINE, line_width=0.5)
    add_text_box(slide, x + 0.15, 4.95, 7.4, 2.6, desc,
                 font_size=9.5, color=COLOR_DARK_TEXT)

# 경쟁사 비교표
add_text_box(slide, 0.8, 7.9, 32, 0.8, "비교 항목별 분석",
             font_size=13, bold=True, color=COLOR_NAVY)

comp_headers = ["비교 항목", "OEM 측 도구", "에임시스템\nNextbe+", "IBM/Siemens", "nori-tc"]
comp_col_w = [7.0, 5.5, 6.0, 6.0, 8.87]
comp_rows = [
    ("TC 단독 도입",       "해당 없음",           "MES 전체 필요",    "TC 분리 불가",     "가능"),
    ("신규 설비 온보딩",   "설비 제조사 탑재",     "엔지니어링 서비스", "MRM 시나리오",    "문서→evidence\n→명세 초안→승인"),
    ("변경 관리",          "해당 없음",           "미확인",           "NPI (고비용)",     "diff+영향도\n+단계반영+복귀"),
    ("AI 기반 운영",       "Exensio 별도 연계",   "공개 자료 없음",   "디지털 트윈/SPC", "온프레미스 LLM\n+RAG 내장"),
    ("도입 비용",          "OEM 라이선스",        "구축+운영 통합",   "엔터프라이즈 고비용", "TC 기능 단독"),
]

x = 0.8
for j, (h, w) in enumerate(zip(comp_headers, comp_col_w)):
    add_rect(slide, x, 8.8, w, 0.85, COLOR_NAVY if j != 4 else COLOR_BLUE)
    add_text_box(slide, x + 0.1, 8.88, w - 0.2, 0.7, h,
                 font_size=9.5, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    x += w

for i, row_data in enumerate(comp_rows):
    x = 0.8
    top = 9.65 + i * 1.4
    bg = COLOR_WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF8)
    for j, (text, w) in enumerate(zip(row_data, comp_col_w)):
        cell_bg = bg if j != 4 else (RGBColor(0xE8, 0xF4, 0xFF) if i % 2 == 0 else RGBColor(0xD0, 0xE8, 0xF8))
        add_rect(slide, x, top, w, 1.35, cell_bg,
                 line_color=COLOR_GRAY_LINE, line_width=0.3)
        tc = COLOR_BLUE if j == 4 else (COLOR_NAVY if j == 0 else COLOR_DARK_TEXT)
        add_text_box(slide, x + 0.1, top + 0.1, w - 0.2, 1.15, text,
                     font_size=9, bold=(j == 4), color=tc, align=PP_ALIGN.CENTER)
        x += w


# ══════════════════════════════════════════════════════════════════
# 슬라이드 12 — 시장 경쟁사 비교 (유형 2: Non-SECS CIM/TC)
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "시장 경쟁사 비교 — Non-SECS CIM/TC 플랫폼형",
                 "nori-tc와 가장 직접적 경쟁 구도 (유형 2)", 12)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 설명
add_rect(slide, 0.8, 2.8, 32, 1.2, COLOR_NAVY)
add_text_box(slide, 1.0, 2.9, 31.5, 1.0,
             "이 유형의 공통 특징: '연결은 되지만 운영은 못 한다' — 설비를 MES에 붙이는 것은 해결하지만 "
             "신규 설비 도입 공수 반복·변경 실패 리스크·장애 추적 불가·AI 부재라는 구조적 문제는 동일하게 남음",
             font_size=10.5, color=COLOR_WHITE)

# 비교표
comp_headers = ["비교 항목", "CIMPlus\nCIMStudio", "AIM\nFMS/CIM", "Inno-CIM", "W-CIM", "PbotEngine", "nori-tc"]
col_w = [6.5, 4.3, 4.3, 4.3, 4.3, 4.3, 5.77]

x = 0.8
for j, (h, w) in enumerate(zip(comp_headers, col_w)):
    hc = COLOR_NAVY if j == 0 else (COLOR_BLUE if j == len(comp_headers) - 1 else RGBColor(0x44, 0x55, 0x66))
    add_rect(slide, x, 4.2, w, 0.85, hc)
    add_text_box(slide, x + 0.1, 4.28, w - 0.2, 0.7, h,
                 font_size=9, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    x += w

rows_data = [
    ("설비 온보딩",   "엔지니어링 공수 필요",       "PLC Map 수작업",          "수작업 인터페이스 개발", "현장 직접 수정",     "Workflow 다이어그램",    "문서→evidence\n→명세초안→승인"),
    ("변경 관리",     "없음",                        "없음",                    "없음",                   "없음",               "없음",                   "diff+영향도\n+단계반영+복귀"),
    ("메시지 신뢰성", "없음",                        "없음",                    "없음",                   "없음",               "없음",                   "partition+DLQ\n+발송이력"),
    ("AI 기능",       "없음",                        "없음",                    "없음",                   "없음",               "없음",                   "온프레미스\nLLM+RAG 내장"),
    ("EDA/FDC",       "구조 별도 설계",              "수집 가능\n(구조화 미확인)", "없음",                 "없음",               "없음",                   "TC 시점부터\n기준 데이터 축"),
    ("주요 시장",     "이차전지/디스플레이\n(대형 팹)", "디스플레이/태양광",    "반도체 부품",            "디스플레이/반도체",  "이차전지\n(SK계열)",     "프로토콜 무관\n전 제조"),
]

for i, row in enumerate(rows_data):
    x = 0.8
    top = 5.05 + i * 2.0
    bg = COLOR_WHITE if i % 2 == 0 else COLOR_GRAY_BG
    for j, (text, w) in enumerate(zip(row, col_w)):
        cell_bg = bg if j != len(col_w) - 1 else (
            RGBColor(0xE8, 0xF4, 0xFF) if i % 2 == 0 else RGBColor(0xD0, 0xE8, 0xF8))
        add_rect(slide, x, top, w, 1.9, cell_bg,
                 line_color=COLOR_GRAY_LINE, line_width=0.3)
        tc = COLOR_NAVY if j == 0 else (COLOR_BLUE if j == len(col_w) - 1 else COLOR_DARK_TEXT)
        add_text_box(slide, x + 0.1, top + 0.15, w - 0.2, 1.6, text,
                     font_size=9, bold=(j == 0 or j == len(col_w) - 1), color=tc,
                     align=PP_ALIGN.CENTER)
        x += w

# 하단 메시지
add_rect(slide, 0.8, 17.1, 32.0, 1.3, COLOR_NAVY)
add_text_box(slide, 1.0, 17.2, 31.5, 1.1,
             "인사이트온 PbotEngine이 nori-tc와 가장 근접한 포지션이지만 Work/Lot/ControlJob 도메인 모델과 변경관리 계층이 없습니다.",
             font_size=11, color=COLOR_WHITE)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 13 — 도입 기대 효과
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "고객사 도입 기대 효과",
                 "신규 설비 도입 · 운영 변경 · 장애 대응 · 장기 고도화", 13)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

effect_sections = [
    ("신규 설비 도입", COLOR_ACCENT,
     ["설비 문서 해석 시간 단축",
      "parser / workflow / 테스트 준비 시간 단축",
      "벤더 확인 질문과 누락 항목 정리 체계화",
      "온보딩 리드타임 30~50% 단축 목표"]),
    ("운영 변경", COLOR_BLUE,
     ["변경 전 검증이 강화되어 사고 확률 감소",
      "사고 발생 시 복귀와 원인 추적이 빨라짐",
      "변경 승인/리뷰 절차의 근거가 명확해짐",
      "긴급 롤백 시간 60~80% 단축 목표"]),
    ("장애 대응", COLOR_GREEN,
     ["MTTR 감소 (사건 타임라인 제공)",
      "유실/중복/반쪽 처리 조사 시간 감소",
      "숙련 운영 인력 의존도 감소",
      "원인 파악 시간 50~70% 단축 목표"]),
    ("장기 고도화", COLOR_PURPLE,
     ["EDA/FDC 연계 비용 감소",
      "데이터 재정의와 중복 수집 비용 감소",
      "상위 시스템 도입 시 기간 단축",
      "EDA/FDC 착수 시 수 개월 단축 가능"]),
]

for i, (title, color, items) in enumerate(effect_sections):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 16.6
    top = 2.9 + row * 7.5
    add_rect(slide, x, top, 16.0, 1.0, color)
    add_text_box(slide, x + 0.2, top + 0.1, 15.6, 0.8, title,
                 font_size=15, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, top + 1.0, 16.0, 6.2, COLOR_WHITE,
             line_color=color, line_width=1.2)
    for j, item in enumerate(items):
        # 체크 아이콘 대체
        add_rect(slide, x + 0.3, top + 1.4 + j * 1.35, 0.5, 0.6, color)
        add_text_box(slide, x + 0.3, top + 1.45 + j * 1.35, 0.5, 0.55, "✓",
                     font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 1.0, top + 1.4 + j * 1.35, 14.7, 0.75, item,
                     font_size=11.5, color=COLOR_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 14 — 도입 방식 및 전환 절차
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "도입 방식 및 기존 시스템 전환 절차",
                 "신규 라인 도입 또는 TC 교체 시점 → 단계적 절체", 14)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

# 두 가지 도입 경로
add_text_box(slide, 0.8, 2.8, 20, 0.8, "주요 도입 경로",
             font_size=13, bold=True, color=COLOR_NAVY)

routes = [
    ("신규 라인 도입", COLOR_BLUE,
     "처음부터 nori-tc 기반으로 설비 연결·모델 정의·운영 절차 구성. "
     "레거시 TC와의 충돌 없이 최적 구성 가능."),
    ("기존 TC 교체", COLOR_ORANGE,
     "4단계 순차 절체 절차를 통해 안전하게 전환. "
     "Spec-to-Proof로 기존 문서 흡수, Verified ChangeOps로 단계 반영."),
]

for i, (title, color, desc) in enumerate(routes):
    x = 0.8 + i * 16.5
    add_rect(slide, x, 3.8, 15.8, 0.85, color)
    add_text_box(slide, x + 0.2, 3.88, 15.4, 0.7, title,
                 font_size=12, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, 4.65, 15.8, 1.7, COLOR_WHITE,
             line_color=color, line_width=1.0)
    add_text_box(slide, x + 0.2, 4.75, 15.4, 1.5, desc,
                 font_size=10.5, color=COLOR_DARK_TEXT)

# 기존 TC 전환 4단계
add_text_box(slide, 0.8, 6.7, 30, 0.8, "기존 TC에서 전환하는 경우 — 4단계 절차",
             font_size=13, bold=True, color=COLOR_NAVY)

transfer_steps = [
    ("STEP 1\n설비 모델\n이관",
     "기존 설비 정의·메시지 규칙·파라미터를\nnori-tc 스키마로 매핑.\nSpec-to-Proof 문서 추출 기능 활용."),
    ("STEP 2\n병행 운영",
     "핵심 설비는 기존 TC와\nnori-tc를 일정 기간 병행 운영하며\n동작 검증."),
    ("STEP 3\n단계적 절체",
     "설비 그룹 단위로 순차 전환.\n전체 라인을 한 번에 전환하지 않음.\nVerified ChangeOps 단계 반영 개념 동일."),
    ("STEP 4\n운영 이관\n완료",
     "검증 완료 설비 그룹부터\n기존 TC 연결 해제 →\nnori-tc 단독 운영."),
]

step_colors = [COLOR_ACCENT, COLOR_BLUE, COLOR_GREEN, COLOR_NAVY]

for i, (title, desc) in enumerate(transfer_steps):
    x = 0.8 + i * 8.1
    # 번호 + 제목
    add_rect(slide, x, 7.7, 7.7, 1.6, step_colors[i])
    add_text_box(slide, x + 0.15, 7.78, 7.4, 1.45, title,
                 font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # 설명
    add_rect(slide, x, 9.3, 7.7, 4.5, COLOR_WHITE,
             line_color=step_colors[i], line_width=1.2)
    add_text_box(slide, x + 0.15, 9.45, 7.4, 4.2, desc,
                 font_size=10.5, color=COLOR_DARK_TEXT)
    # 화살표
    if i < len(transfer_steps) - 1:
        add_text_box(slide, x + 7.7, 8.2, 0.35, 0.8, "→",
                     font_size=20, bold=True, color=COLOR_NAVY, align=PP_ALIGN.CENTER)

# 하단 메시지
add_rect(slide, 0.8, 14.2, 32.0, 1.5, RGBColor(0xE8, 0xF4, 0xE8))
add_rect(slide, 0.8, 14.2, 0.4, 1.5, COLOR_GREEN)
add_text_box(slide, 1.4, 14.3, 31.0, 1.3,
             "전환 과정에서 기존 설비 문서와 운영 이력은 nori-tc 내부의 기준 데이터로 흡수되어,\n"
             "이후 AI 기능의 RAG 기반 지식으로도 활용됩니다.",
             font_size=11, color=COLOR_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 15 — 결론
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "결론 — nori-tc 제안 포인트",
                 "라인 실행 리스크를 낮추고 생산 고도화의 기반을 만드는 실행 플랫폼", 15)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_NAVY)
add_rect(slide, 0, 2.5, 0.5, 16.2, COLOR_ACCENT)

# 6개 핵심 포인트
conclusion_items = [
    ("Spec-to-Proof",                    COLOR_ACCENT,
     "원본 설비 문서를 검토 가능한 명세로 더 빠르게 전환"),
    ("Verified ChangeOps",               COLOR_BLUE,
     "운영 중 변경을 더 안전하게 다룸"),
    ("Provable Reliability",             COLOR_LIGHT_BLUE,
     "메시지와 처리 결과를 근거 기반으로 설명 가능하게"),
    ("Explainable Operation Intelligence", COLOR_GREEN,
     "장애 대응을 사건 단위로 더 쉽게"),
    ("Industrial Data Spine",            COLOR_PURPLE,
     "EDA/FDC/AI 활용으로 자연스럽게 이어지는 데이터 축"),
    ("Autonomous Commissioning Lab",     COLOR_ORANGE,
     "실라인 반영 전 동일한 엔진으로 사전 검증"),
]

for i, (kw, color, desc) in enumerate(conclusion_items):
    col = i % 2
    row = i // 2
    x = 1.0 + col * 16.4
    top = 3.2 + row * 4.2
    add_rect(slide, x, top, 15.5, 1.1, color)
    add_text_box(slide, x + 0.2, top + 0.1, 15.1, 0.9, kw,
                 font_size=14, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, top + 1.1, 15.5, 2.7, RGBColor(0x1E, 0x44, 0x72))
    add_text_box(slide, x + 0.2, top + 1.3, 15.1, 2.3, desc,
                 font_size=12, color=RGBColor(0xBB, 0xCC, 0xDD))

# 최종 메시지
add_rect(slide, 0.8, 16.0, 32.0, 2.3, COLOR_ACCENT)
add_text_box(slide, 1.0, 16.15, 31.5, 2.0,
             "nori-tc는 단순 연결 시스템이 아니라,\n라인 실행 리스크를 낮추고 이후 생산 고도화의 기반을 만드는 실행 플랫폼입니다.",
             font_name=FONT_TITLE, font_size=15, bold=True,
             color=COLOR_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# 슬라이드 16 — 기술 실사 (마지막)
# ══════════════════════════════════════════════════════════════════
slide = add_new_slide()
add_slide_header(slide, "기술 실사 (Technical Due Diligence)",
                 "현재 기반 항목은 실제 구현 코드·DB 스키마·운영 설계 문서를 기반으로 합니다", 16)

add_rect(slide, 0, 2.5, 33.87, 16.2, COLOR_GRAY_BG)

add_text_box(slide, 0.8, 2.8, 32, 0.9,
             "도입 검토 단계에서 기술 실사가 필요하시면 아래 항목에 대한 상세 자료를 별도로 제공드릴 수 있습니다.",
             font_size=12, color=COLOR_DARK_TEXT)

due_items = [
    ("Gateway 파티션 라우팅",    "route_partition 기반 라우팅 규칙 설계 문서",           COLOR_BLUE),
    ("메시지 처리 흐름",          "메시지 처리 흐름 및 DB 스키마 정의",                   COLOR_NAVY),
    ("작업 계층 구조",            "Work→ControlJob→ProcessJob→Lot 도메인 모델 정의",       COLOR_GREEN),
    ("Dual Response",            "Dual Response 수집 및 Async 결과 저장 구조 설명",        COLOR_LIGHT_BLUE),
    ("설비 파라미터 관리",        "Check Out / Check In 운영 절차 문서",                   COLOR_PURPLE),
    ("모델 버전 관리",            "모델 버전 관리 API 명세",                               COLOR_ORANGE),
]

for i, (title, desc, color) in enumerate(due_items):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 16.5
    top = 4.0 + row * 4.2
    add_rect(slide, x, top, 16.0, 0.9, color)
    add_text_box(slide, x + 0.2, top + 0.1, 15.6, 0.75, title,
                 font_size=13, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, top + 0.9, 16.0, 2.8, COLOR_WHITE,
             line_color=color, line_width=1.0)
    add_text_box(slide, x + 0.2, top + 1.1, 15.6, 2.4, desc,
                 font_size=11, color=COLOR_DARK_TEXT)

# 하단 연락처 영역
add_rect(slide, 0.8, 16.5, 32.0, 1.8, COLOR_NAVY)
add_text_box(slide, 1.0, 16.65, 31.5, 1.5,
             "기술 실사 문의 및 POC(Proof of Concept) 진행을 원하시면 연락주시기 바랍니다.\n"
             "nori-tc 개발팀 · 실제 코드·스키마·설계 문서를 기반으로 상세 설명 드립니다.",
             font_size=11, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


# ─── 저장 ───
import os
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"저장 완료: {OUTPUT_PATH}")
print(f"총 슬라이드 수: {len(prs.slides)}")
EOF
